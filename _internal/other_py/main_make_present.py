import comtypes.client
import time
import main_win as main
import pptx as pp
import os
import _internal.other_py.other_funcs as oth
from tkinter import messagebox as msgbox
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE_TYPE
from copy import deepcopy
from io import BytesIO
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


def copy_shape_with_crop(source_shape, target_slide):
    """
    Копирует фигуру (включая картинки с кропом) через прямое копирование XML.
    Обходит баг python-pptx с внутренними методами упаковки изображений.
    """
    # 1. Глубокая копия XML-элемента фигуры
    new_el = deepcopy(source_shape.element)

    if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        blip = new_el.find('.//a:blip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
        if blip is not None and '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed' in blip.attrib:
            # Получаем байты исходного изображения
            old_rid = blip.get(qn('r:embed'))
            image_part = source_shape.part.related_part(old_rid)

            # 🔑 Безопасная регистрация изображения через публичный API
            temp_pic = target_slide.shapes.add_picture(
                BytesIO(image_part.blob), 0, 0, 1, 1  # временные координаты
            )

            # Извлекаем корректный rId из временной фигуры
            temp_blip = temp_pic._element.find('.//a:blip', namespaces={
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            new_rid = temp_blip.get(qn('r:embed'))

            # Удаляем временную фигуру из XML (само изображение остаётся в пакете .zip)
            temp_pic._element.getparent().remove(temp_pic._element)

            # Подставляем новый rId в скопированную фигуру
            blip.set(qn('r:embed'), new_rid)

    # 2. Генерация нового уникального ID для фигуры
    cnvpr = new_el.find('.//p:cNvPr', namespaces={'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
    if cnvpr is not None:
        max_id = max((s.shape_id for s in target_slide.shapes), default=0)
        cnvpr.set('id', str(max_id + 1))

    # 3. Вставка элемента в дерево слайда
    target_slide.shapes._spTree.append(new_el)


def main_make_present(old_path_pptx, exel_path, file_type, dir_name, main_name_for_file, new_file_name, inputs, call = None):

    data, names_list = oth.corr_exel_data(exel_path, inputs, main_name_for_file)

    if data:
        old_prs = pp.Presentation(old_path_pptx)
        new_prs = pp.Presentation()

        orig_width = old_prs.slide_width
        orig_height = old_prs.slide_height

        new_prs.slide_width = orig_width
        new_prs.slide_height = orig_height

        maket_old_slide = old_prs.slides[0]
        layout_new_slide = new_prs.slide_layouts[6]

        for i in range(len(data)):
            new_slide = new_prs.slides.add_slide(layout_new_slide)

            # Копируем все фигуры из шаблона с сохранением кропа
            for shape in maket_old_slide.shapes:
                copy_shape_with_crop(shape, new_slide)

            # Заменяем текст в фигурах
            for shape in new_slide.shapes:
                if not shape.has_text_frame:
                    continue

                text = shape.text.strip()
                if not text:
                    continue

                if text in data[0]:
                    index_word = data[0].index(text)
                    oth.replace_text_keep_format(shape, data[i][index_word])

        os.makedirs(os.path.split(old_path_pptx)[0] + '/' + dir_name, exist_ok=True)
        path = os.path.split(old_path_pptx)[0] + '/' + dir_name + '/' + str(new_file_name) + '.pptx'
        new_prs.save(path)

        base_dir = os.path.dirname(os.path.abspath(old_path_pptx))
        out_dir = os.path.join(base_dir, dir_name)
        os.makedirs(out_dir, exist_ok=True)

        abs_pptx_path = os.path.normpath(os.path.abspath(path))



        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = True

        try:


            deck = powerpoint.Presentations.Open(abs_pptx_path)


            time.sleep(0.5)


            for i in range(1, deck.Slides.Count + 1):
                slide = deck.Slides(i)


                filename = os.path.join(out_dir, f"{names_list[i-1]}_{i}.{file_type}")

                abs_filename = os.path.normpath(os.path.abspath(filename))

                try:

                    slide.Export(abs_filename, file_type)

                except Exception as e:
                    pass

        finally:

            try:
                deck.Close()
            except:
                pass

            try:
                powerpoint.Quit()
            except:
                pass


            del deck
            del powerpoint
            call()


    else:
        msgbox.showerror(message='Колонки для названия файла не нашлось')

        main.main()

if __name__ == "__main__":
    main_make_present("D:/My things/Projects/Проект АвтоГрамоты/Шаблон 1.pptx", "D:/My things/Projects/Проект АвтоГрамоты/Таблица-шаблоны.xlsx",'PNG','privet','Имя','hoyhoy', {'ОО':'Образовательная организация'})

