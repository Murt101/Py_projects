import io, openpyxl
import tkinter as tk
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE
from _internal.other_py import main_make_present as mmp
import threading

stop_ender_page = [False,False, "dir_name"]


def corr_exel_data(exel_path=None, inp: dict = None, ie = None):
    list_manes = []
    exel_s = openpyxl.load_workbook(exel_path, read_only=True)
    exel = exel_s.active

    data = [[str(cell) if cell is not None else "" for cell in row]
            for row in exel.iter_rows(values_only=True)]
    exel_s.close()

    if ie in data[0]:
        index_name = data[0].index(ie)
        for j in range(len(data)):
            list_manes.append(data[j][index_name])
        if inp == 'AUTO':
            pass
        else:
            val_to_key = {v: k for k, v in inp.items()}
            for row in data:
                for i, item in enumerate(row):
                    if item in val_to_key:
                        row[i] = val_to_key[item]


        return data, list_manes
    else:
        return False, False

def copy_shape_to_slide(source_shape, target_slide):
    shape_type = source_shape.shape_type

    # 1. КАРТИНКИ
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            # Оборачиваем байты в файлоподобный объект
            image_stream = io.BytesIO(source_shape.image.blob)
            new_shape = target_slide.shapes.add_picture(
                image_stream,
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )
            if hasattr(source_shape, 'rotation') and source_shape.rotation is not None:
                new_shape.rotation = source_shape.rotation
        except Exception as e:
            pass

    # 2. АВТОФИГУРЫ
    elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        new_shape = target_slide.shapes.add_shape(
            source_shape.auto_shape_type,
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )
        if source_shape.has_text_frame:
            _copy_text_frame(source_shape.text_frame, new_shape.text_frame)
        _copy_fill_style(source_shape, new_shape)
        _copy_line_style(source_shape, new_shape)
        if hasattr(source_shape, 'rotation') and source_shape.rotation is not None:
            new_shape.rotation = source_shape.rotation

    # 3. ТЕКСТОВЫЕ ПОЛЯ
    elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        new_shape = target_slide.shapes.add_textbox(
            source_shape.left,
            source_shape.top,
            source_shape.width,
            source_shape.height
        )
        if source_shape.has_text_frame:
            _copy_text_frame(source_shape.text_frame, new_shape.text_frame)
        _copy_fill_style(source_shape, new_shape)
        _copy_line_style(source_shape, new_shape)

    else:
        pass  # Таблицы, диаграммы, группы пропускаем

def _copy_text_frame(source_tf, target_tf):
    if not source_tf.paragraphs: return
    target_tf.clear()

    for i, para in enumerate(source_tf.paragraphs):
        target_para = target_tf.paragraphs[0] if i == 0 else target_tf.add_paragraph()
        target_para.alignment = para.alignment
        target_para.space_before = para.space_before
        target_para.space_after = para.space_after
        if para.line_spacing is not None:
            target_para.line_spacing = para.line_spacing

        for run in para.runs:
            new_run = target_para.add_run()
            new_run.text = run.text
            f_src, f_tgt = run.font, new_run.font

            if f_src.bold is not None: f_tgt.bold = f_src.bold
            if f_src.italic is not None: f_tgt.italic = f_src.italic
            if f_src.underline is not None: f_tgt.underline = f_src.underline
            if f_src.size is not None: f_tgt.size = f_src.size
            if f_src.name is not None: f_tgt.name = f_src.name

            if f_src.color.type == MSO_COLOR_TYPE.RGB:
                f_tgt.color.rgb = f_src.color.rgb

def _copy_fill_style(source_shape, target_shape):
    """Безопасное копирование заливки с защитой от _NoFill и градиентов"""
    src, tgt = source_shape.fill, target_shape.fill
    try:
        # Проверяем сплошную заливку и тип цвета
        if src.solid and src.fore_color.type == MSO_COLOR_TYPE.RGB:
            tgt.solid()
            tgt.fore_color.rgb = src.fore_color.rgb
        else:
            tgt.background()
    except (TypeError, AttributeError):
        # Если заливка отсутствует, градиентная, тематическая или битая -> прозрачная
        tgt.background()

def _copy_line_style(source_shape, target_shape):
    """Безопасное копирование контура"""
    src, tgt = source_shape.line, target_shape.line
    try:
        if src.width is not None:
            tgt.width = src.width
        if hasattr(src, 'dash_style') and src.dash_style is not None:
            tgt.dash_style = src.dash_style
        if src.color.type == MSO_COLOR_TYPE.RGB:
            tgt.color.rgb = src.color.rgb
        elif src.color.type == MSO_COLOR_TYPE.AUTO:
            tgt.color.auto()
    except Exception:
        pass

def replace_text_keep_format(shape, new_text):
    """Заменяет текст, сохраняя шрифт, размер, цвет и жирность из шаблона"""
    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    # Исправлено: было parframes, стало paragraphs
    if not tf.paragraphs:
        return

    p = tf.paragraphs[0]

    # Удаляем все текстовые фрагменты (runs) кроме первого
    for i in range(len(p.runs) - 1, 0, -1):
        p._element.remove(p._element[i])

    # Вставляем новый текст в первый run (сохраняя его стиль)
    if p.runs:
        p.runs[0].text = new_text
    else:
        p.add_run().text = new_text

def ender_page____(old_path_pptx, exel_path, file_type, dir_name, main_name_for_file, new_file_name, inputs):

    def callback__():
        main_text['text'] = 'hi'

    ender_win = tk.Tk()
    ender_win.geometry('300x300')
    ender_win.resizable(False, False)

    main_text = tk.Button(ender_win, text = 'hellow')
    main_text.pack()

    threading.Thread(
        target=mmp.main_make_present,
        kwargs={
            "old_path_pptx": old_path_pptx,
            "exel_path": exel_path,
            "file_type": file_type,
            "dir_name": dir_name,
            "main_name_for_file": main_name_for_file,
            "new_file_name": new_file_name,
            "inputs": inputs,
            'call': callback__
        },
        daemon=True
    ).start()



    ender_win.mainloop()

if '__main__' == __name__:
    replace_text_keep_format('','')
    copy_shape_to_slide('','')
    _copy_text_frame('','')
    _copy_fill_style('','')
    _copy_line_style('','')
    corr_exel_data()
